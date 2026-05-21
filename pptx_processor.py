import os
import logging
from pptx import Presentation
from typing import List
from ner_engine import extract_all_entities
from docx_processor import chunk_text

logger = logging.getLogger("docguard")

def process_pptx(input_path: str, output_path: str):
    logger.info(f"  ▶ 正在处理 PPT 文档: {os.path.basename(input_path)}")
    try:
        prs = Presentation(input_path)
    except Exception as e:
        logger.error(f"  ❌ 读取 {input_path} 失败: {e}")
        return False, ""
        
    # 1. 提取所有幻灯片中的文本
    full_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                full_text.append(shape.text)
                
    raw_text = "\n".join(full_text)
    if not raw_text.strip():
        logger.warning(f"  [警告] {os.path.basename(input_path)} 未提取到文本。")
        return False, ""
        
    # 2. 分块送给大模型进行实体识别
    text_chunks = chunk_text(raw_text)
    entities = extract_all_entities(text_chunks)
    logger.info(f"    ✅ 共找到 {len(entities)} 个敏感实体: {entities[:5]}...")
    
    # 3. 遍历文本框中的 Runs，执行精确替换 (保持样式)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text:
                        for entity in entities:
                            if entity in run.text:
                                run.text = run.text.replace(entity, "[脱敏]")
                                
    # 4. 保存
    prs.save(output_path)
    
    clean_text = raw_text
    for entity in entities:
        clean_text = clean_text.replace(entity, "[脱敏]")
        
    return True, clean_text
